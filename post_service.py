from flask import Flask, request, jsonify, send_file
import pandas as pd
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
import numpy as np

app = Flask(__name__)

# Load data once when the app starts
sheet_name = 'Base year_2021'
usecols = 'B:AK'
df = pd.read_excel('20230411_A1OI_Summary Update_v4.xlsx', sheet_name=sheet_name, usecols=usecols)
priority_companies_df = pd.read_excel('A1OI Priority client mapping.xlsx')
priority_companies = priority_companies_df['Informal client name'].tolist()
df['Priority_status'] = np.where(df['Company'].isin(priority_companies_df['Informal client name']), 'Priority', 'Non-Priority')
df['Geography'] = df['Geography'].replace('N. America', 'AMER')
df = df.fillna(0)

def filter_data(df, industry, relationship, revenue, geography, country, industry_primary, combined_opp_score, priority_status):
    df['Country'] = df['Country'].str.strip()

    if isinstance(country, str):
        country = [c.strip() for c in country.split(',')]
    else:
        country = [c.strip() for c in country]

    if industry:
        df = df[df['Bain Industry'] == industry]
    if industry_primary:
        df = df[df['Primary Industry'] == industry_primary]
    if relationship:
        if relationship == 'Current':
            df = df[df['Bain Relationship'].str.startswith('Current')]
        elif relationship == 'Recent':
            df = df[df['Bain Relationship'].str.startswith('Recent')]
        elif relationship == 'Greyspace':
            df = df[df['Bain Relationship'].str.startswith('Greyspace')]
    if revenue == 'less_than_500000':
        df = df[df['Revenue (M)'] < 5000]
    elif revenue == 'greater_than_500000':
        df = df[df['Revenue (M)'] > 5000]
    if priority_status:
        df = df[df['Priority_status'] == priority_status]
    if geography:
        df = df[df['Geography'] == geography]
    if country:
        df = df[df['Country'].isin(country)]
    if combined_opp_score:
        if combined_opp_score == 'less_than_5':
            df = df[df['Combined Opp. Score'] < 5]
        elif combined_opp_score == 'between_5_and_10':
            df = df[(df['Combined Opp. Score'] >= 5) & (df['Combined Opp. Score'] <= 10)]
        elif combined_opp_score == 'greater_than_10':
            df = df[df['Combined Opp. Score'] > 10]
    
    return df

@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    industry = request.form.get('industry')
    relationship = request.form.get('relationship')
    revenue = request.form.get('revenue')
    country = request.form.get('country[]', '').split('|')
    country = [c.strip() for c in country if c.strip()]
    geography = request.form.get('geography')
    industry_primary = request.form.get('industry_primary')
    combined_opp_score = request.form.get('combined_opp_score')
    priority_status = request.form.get('priority_status')

    df_filtered = filter_data(df, industry, relationship, revenue, geography, country, industry_primary, combined_opp_score, priority_status)

    ppt_path = generate_ppt_file(df_filtered, industry, country, revenue, geography, industry_primary, combined_opp_score, priority_status)

    return send_file(ppt_path, as_attachment=True)

def generate_ppt_file(df_filtered, industry, country, revenue, geography, industry_primary, combined_opp_score, priority_status):
    script_dir = os.path.dirname(os.path.abspath(__file__))
    file_path = os.path.join(script_dir, 'Python to PPT', 'inputs', 'Template_v1.pptx')
    country_str = ', '.join(country) if country else 'No Country Selected'
    file_path_1 = os.path.join(script_dir, 'Python to PPT', 'Outputs', f'Template_tables_{industry}_{country_str}.pptx')

    presentation = Presentation(file_path)

    def populate_table_with_data(table, df, font_size):
        center_align_columns = [1, 4, 6, 7, 8, 9, 10]
        comma_separator_columns = [6, 7, 8, 9]

        for i in range(1, min(len(df) + 1, len(table.rows))):
            for j in range(min(len(df.columns), len(table.columns))):
                cell = table.cell(i, j)
                value = df.iloc[i - 1, j]
                if isinstance(value, str) and value.strip() == '':
                    cell_value = ''
                else:
                    if j == 1:
                        cell_value = round(float(value) / 1000, 1)
                    elif j == 4:
                        cell_value = round(float(value) * 100, 1)
                    elif j == 10:
                        cell_value = f"{value}x" if isinstance(value, float) else value
                    elif j in comma_separator_columns:
                        cell_value = f"{int(float(value)):,}"
                    else:
                        cell_value = value

                cell.text = str(cell_value)

                for paragraph in cell.text_frame.paragraphs:
                    pPr = paragraph._element.get_or_add_pPr()
                    buNone = OxmlElement("a:buNone")
                    pPr.insert(0, buNone)
                    paragraph.space_before = Pt(0)
                    paragraph.space_after = Pt(0)
                    paragraph.line_spacing = 1.0
                    paragraph.level = 0

                    if j in center_align_columns:
                        paragraph.alignment = PP_ALIGN.CENTER

                    for run in paragraph.runs:
                        run.font.size = font_size
                        run.font.color.rgb = RGBColor(0, 0, 0)
                        if j == 0:
                            run.font.bold = False

                    ind = pPr.find(".//a:ind", namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                    if ind is not None:
                        pPr.remove(ind)

                    new_ind = OxmlElement('a:ind')
                    new_ind.set('left', '0')
                    new_ind.set('hanging', '0')
                    pPr.append(new_ind)

                lstStyle = cell.text_frame._element.find(".//a:lstStyle", namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                if lstStyle is not None:
                    cell.text_frame._element.remove(lstStyle)

    relationship_types = ['Current', 'Recent', 'Greyspace']
    selected_columns = [
        'Company','Revenue (M)','Country', 'Primary Industry', 'EBIT margin Change (% pts.)  [Benchmark year-Base year]',
        'Relative-EBIT Margin','Static EBIT opportunity (Median)', 'Static SG&A Opportunity (Median)', 
        'Static NWC opportunity (Median)','Median Incremental Rev.- EBIT Oppt. (Static Historical)',
        "Net debt/EBITDA (Dec'22)", 'Relative ESG Combined score'
    ]

    slide_index = 1
    for relationship in relationship_types:
        filtered_df = filter_data(df_filtered, industry, relationship, revenue, geography, country, industry_primary, combined_opp_score, priority_status)
        filtered_df = filtered_df[selected_columns]
        print(filtered_df)
        if not filtered_df.empty:
            num_slides = (len(filtered_df) + 9) // 10
            for slide_num in range(num_slides):
                df_slice = filtered_df[slide_num * 10:(slide_num + 1) * 10]
                slide = presentation.slides[slide_index]
                title_shape = slide.shapes.title
                revenue_str = '<$5B' if revenue == 'less_than_500000' else '>$5B' if revenue == 'greater_than_500000' else "All"
                title_shape.text = f"{industry} ({country_str}) Revenue: {revenue_str} {priority_status} ({relationship}) - ({slide_num + 1}/{num_slides})"

                table = next((shape.table for shape in slide.shapes if shape.has_table), None)
                if table:
                    populate_table_with_data(table, df_slice, font_size=Pt(10))

                slide_index += 1

    while slide_index < len(presentation.slides):
        slide = presentation.slides[slide_index]
        xml_slides = presentation.slides._sldIdLst
        slide_id = slide.slide_id
        for s in xml_slides:
            if s.attrib['id'] == str(slide_id):
                xml_slides.remove(s)
                break

    presentation.save(file_path_1)

    return file_path_1

if __name__ == '__main__':
    app.run(debug=True, port=5001)
