from flask import Flask, render_template, request, send_file, redirect, url_for, flash
import pandas as pd
import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.chart.data import BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
import numpy as np


app = Flask(__name__)

app.secret_key = 'supersecretkey'  # Needed for flashing messages
# Define the correct access key
ACCESS_KEY = "wkdqnbrxylrudkxoorkdzkzlq"

# Define the folder to store uploaded files
INPUT_FOLDER = 'input_folder/'
app.config['INPUT_FOLDER'] = INPUT_FOLDER

# Ensure the input folder exists
if not os.path.exists(INPUT_FOLDER):
    os.makedirs(INPUT_FOLDER)

# Load data once when the app starts
filename1 = 'PI data.xlsx'
filename2 = 'Client mapping.xlsx'
sheet_name = 'Base year_2021'
usecols = 'B:AM'
# Load the main dataset and priority companies data
script_dir = os.path.dirname(os.path.abspath(__file__))
input_path = os.path.join(script_dir, 'input_folder', filename1)
# print(script_dir)
df = pd.read_excel(input_path, sheet_name=sheet_name, usecols=usecols)
mapping_path = os.path.join(script_dir, 'input_folder', filename2)
priority_companies_df = pd.read_excel(mapping_path)
priority_companies = priority_companies_df['Informal client name'].tolist()
# Add a column to indicate priority status based on the list of priority companies
df['Priority_status'] = np.where(df['Company'].isin(priority_companies_df['Informal client name']), 'Priority', 'Non-Priority')
df['Primary Industry (Sector)'] = df['Primary Industry (Sector)'].astype(str)
# Replace certain geographic values and fill NaN values with 0
df['Geography'] = df['Geography'].replace('N. America', 'AMER')
df = df.fillna(0)

# Function to get unique values from the dataset for filtering
def get_contextual_values(df,filtered_data, geography=None, industry=None):
    if geography:
        filtered_df = df[df['Geography'] == geography]
    else:
        filtered_df=df

    if industry:
        filtered_df = filtered_df[filtered_df['Bain Industry'] == industry]
      
        
    
    geographies = sorted(df['Geography'].unique().tolist())
    countries = sorted(filtered_data['Country'].unique().tolist()) if geography else sorted(df['Country'].unique().tolist())
    industries = sorted(df['Bain Industry'].unique().tolist())
    industries_primary = sorted(filtered_data['Primary Industry'].unique().tolist())
    industries_primary_sector = sorted(filtered_data['Primary Industry (Sector)'].unique().tolist())
    relationships = sorted(df['Bain Relationship (updated)'].unique().tolist())
    priority_column = sorted(df['Priority_status'].unique().tolist())
    companies = sorted(filtered_data['Company'].unique().tolist())

    return industries, relationships, countries, geographies, industries_primary, priority_column, industries_primary_sector, companies

# Function to filter the dataset based on various criteria
def filter_data(df, industry, relationship, revenue, geography, country, industry_primary, combined_opp_score, priority_status, industry_primary_sector, company):
    df['Company'] = df['Company'].str.strip()
    df['Country'] = df['Country'].str.strip()
    df['Bain Relationship (updated)'] = df['Bain Relationship (updated)'].str.strip()
    df['Primary Industry'] = df['Primary Industry'].str.strip()



    if industry == "Select Industry" and industry_primary == "Select Primary Industry" and relationship == "All relationships" and revenue == "Select Revenue Range" and combined_opp_score == "Select Score Range" and priority_status == "Select Priority" and geography == "Select Region" and country == "Select Country" and company == "All companies":
        return df

    if isinstance(country, str):
        country = [c.strip() for c in country.split(',')]
    else:
        country = [c.strip() for c in country]

    if isinstance(company, str):
        company = [c.strip() for c in company.split(',')]
    else:
        company = [c.strip() for c in company]

    if industry and industry != "Select Industry":
        df = df[df['Bain Industry'] == industry]


    if isinstance(industry_primary, str):
        industry_primary=[c.strip() for c in industry_primary.split(',')]
    else:
        industry_primary = [c.strip() for c in industry_primary]

    if industry_primary_sector:
        df = df[df['Primary Industry (Sector)'] == industry_primary_sector]
    if isinstance(relationship, str):
        relationship=[c.strip() for c in relationship.split(',')]
    else:
        relationship = [c.strip() for c in relationship]
    if revenue and revenue != "Select Revenue Range":
        if revenue == 'less_than_100000':
            df = df[df['Revenue (M)'] < 1000]
        elif revenue == 'between 10000_and_50000':
            df = df[(df['Revenue (M)'] > 1000) & (df['Revenue (M)'] < 5000)]
        elif revenue == 'between 50000_and_100000':
            df = df[(df['Revenue (M)'] > 5000) & (df['Revenue (M)'] < 10000)]
        elif revenue == 'between 100000_and_200000':
            df = df[(df['Revenue (M)'] > 10000) & (df['Revenue (M)'] < 20000)]
        elif revenue == 'more_than_200000':
            df = df[df['Revenue (M)'] > 20000]
    if priority_status and priority_status != "Select Priority":
        df = df[df['Priority_status'] == priority_status]
    if geography and geography != "Select Region":
        df = df[df['Geography'] == geography]
    if relationship and relationship != "All relationships":
        df = df[df['Bain Relationship (updated)'].isin(relationship)]
    if country and country != "Select Country":
        df = df[df['Country'].isin(country)]
    if company and company != "All companies":
        df = df[df['Company'].isin(company)]
    if industry_primary and industry_primary != "Select Primary Industry":
        df = df[df['Primary Industry'].isin(industry_primary)]
    if combined_opp_score and combined_opp_score != "Select Score Range":
        if combined_opp_score == 'less_than_5':
            df = df[df['Combined Opp. Score'] < 5]
        elif combined_opp_score == 'between_5_and_10':
            df = df[(df['Combined Opp. Score'] >= 5) & (df['Combined Opp. Score'] <= 10)]
        elif combined_opp_score == 'greater_than_10':
            df = df[df['Combined Opp. Score'] > 10]
    return df
    


# Route to render the home page
@app.route('/')
def home():
    # Get filter parameters from the request
    industry = request.args.get('industry')
    relationship = request.args.getlist('relationship[]')
    revenue = request.args.get('revenue')
    country = request.args.getlist('country[]')
    company = request.args.getlist('company[]')
    geography = request.args.get('geography')
    industry_primary = request.args.getlist('industry_primary[]')
    industry_primary_sector = request.args.get('industry_primary_sector')
    combined_opp_score = request.args.get('combined_opp_score')
    priority_status = request.args.get('priority_status')

    # Debugging to check if the values are being parsed correctly
    print("Country:", country)
    print("Company:", company)
    print("Relationship", relationship)

    # Check if any filters have been applied
    filters_applied = any([industry, relationship, revenue, country, geography, industry_primary, combined_opp_score, priority_status, industry_primary_sector, company])

    # Filter the dataset based on the selected filters
    df_filtered = filter_data(df, industry, relationship, revenue, geography, country, industry_primary, combined_opp_score, priority_status, industry_primary_sector, company)
    print(df_filtered)
    df_filtered.reset_index(drop=True, inplace=True)
    filtered_data = df_filtered.to_dict(orient='records')
    min_revenue = df_filtered['Revenue (M)'].min()

    # Get contextual values for dropdowns and filters
    industries, relationships, countries, geographies, industries_primary, priority_column, industries_primary_sector, companies = get_contextual_values(df,df_filtered, geography, industry)

    # Render the template with the filtered data and contextual values
    return render_template('index.html',
                           industries=industries,
                           industries_primary=industries_primary,
                           industries_primary_sector=industries_primary_sector,
                           relationships=relationships,
                           geographies=geographies,
                           countries=countries,
                           companies = companies,
                           revenue_ranges=['less_than_100000', 'between 10000_and_50000','between 50000_and_100000','between 100000_and_200000','more_than_200000'],
                           filtered_data=filtered_data,
                           filters_applied=filters_applied,
                           industry=industry,
                           industry_primary=industry_primary,
                           industry_primary_sector=industry_primary_sector,
                           revenue=revenue,
                           relationship=relationship,
                           geography=geography,
                           country=country,
                           company = company,
                           min_revenue=min_revenue,
                           priority_companies=priority_companies,
                           priority_column=priority_column,
                           priority_status=priority_status)


# Route to generate a PowerPoint presentation
@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():

    # Get filter parameters from the form submission
    industry = request.form.get('industry')
    relationship = request.form.get('relationship')
    revenue = request.form.get('revenue')
    country = request.form.get('country[]', '').split('|')
    country = [c.strip() for c in country if c.strip()]
    company = request.form.get('company[]', '').split('|')
    company = [c.strip() for c in company if c.strip()]
    geography = request.form.get('geography')
    industry_primary = request.form.get('industry_primary')
    industry_primary_sector = request.form.get('industry_primary_sector')
    combined_opp_score = request.form.get('combined_opp_score')
    priority_status = request.form.get('priority_status')

        # Function to filter the dataset based on various criteria
    def filter_data_2(df, industry, revenue, geography, country, industry_primary, combined_opp_score, priority_status, industry_primary_sector, company):
        df['Company'] = df['Company'].str.strip()
        df['Country'] = df['Country'].str.strip()
        
        if industry=="Select Industry" and industry_primary=="Select primary Industry" and relationship=="All relationships" and revenue=="Select Revenue Range" and combined_opp_score=="Select Score Range" and priority_status == "Select Priority" and geography=="Select Region" and country == "Select Country" and company == "All companies":
            return df
        if isinstance(country, str):
            country = [c.strip() for c in country.split(',')]
        else:
            country = [c.strip() for c in country]

        if isinstance(company, str):
            company = [c.strip() for c in company.split(',')]
        else:
            company = [c.strip() for c in company]

        if industry:
            df = df[df['Bain Industry'] == industry]
        if industry_primary:
            df = df[df['Primary Industry'] == industry_primary]
        if industry_primary_sector:
            df = df[df['Primary Industry (Sector)'] == industry_primary_sector]
        if relationship:
            if relationship == 'Current':
                df = df[df['Bain Relationship'].str.startswith('Current')]
            elif relationship == 'Recent':
                df = df[df['Bain Relationship'].str.startswith('Recent')]
            elif relationship == 'Greyspace':
                df = df[df['Bain Relationship'].str.startswith('Greyspace')]
            elif relationship == 'Whitespace':
                df = df[df['Bain Relationship'].str.startswith('Whitespace')]
        if revenue == 'less_than_100000':
            df = df[df['Revenue (M)'] < 1000]
        elif revenue == 'between 10000_and_50000':
            df = df[(df['Revenue (M)'] > 1000) & (df['Revenue (M)'] < 5000)]
        elif revenue == 'between 50000_and_100000':
            df = df[(df['Revenue (M)'] > 5000) & (df['Revenue (M)'] < 10000)]
        elif revenue == 'between 100000_and_200000':
            df = df[(df['Revenue (M)'] > 10000) & (df['Revenue (M)'] < 20000)]
        elif revenue == 'more_than_200000':
            df = df[df['Revenue (M)'] > 20000]
        if priority_status:
            df = df[df['Priority_status'] == priority_status]
        if geography:
            df = df[df['Geography'] == geography]
        if country:
            df = df[df['Country'].isin(country)]
        if company:
            df = df[df['Company'].isin(company)]
        if combined_opp_score:
            if combined_opp_score == 'less_than_5':
                df = df[df['Combined Opp. Score'] < 5]
            elif combined_opp_score == 'between_5_and_10':
                df = df[(df['Combined Opp. Score'] >= 5) & (df['Combined Opp. Score'] <= 10)]
            elif combined_opp_score == 'greater_than_10':
                df = df[df['Combined Opp. Score'] > 10]
        return df


    # Filter the dataset based on the selected filters
    df_filtered = filter_data_2(df, industry, revenue, geography, country, industry_primary, combined_opp_score, priority_status, industry_primary_sector, company)
    # filtered_df = df_filtered.to_dict(orient='records')
    # print(df_filtered)
    # Generate a PowerPoint file with the filtered data
    ppt_path = generate_ppt_file(df_filtered, industry, country, revenue, geography, industry_primary, combined_opp_score, priority_status,industry_primary_sector, company)

    # Send the generated PowerPoint file as an attachment
    return send_file(ppt_path, as_attachment=True)


# Function to generate the PowerPoint file based on the filtered data
def generate_ppt_file(df_filtered, industry, country, revenue, geography, industry_primary, combined_opp_score, priority_status,industry_primary_sector, company):
    script_dir = os.path.dirname(os.path.abspath(__file__))
    file_path = os.path.join(script_dir, 'Python to PPT', 'inputs', 'Template_v1.pptx')
    country_str = ', '.join(country) if country else 'All countries'
    output_pptx_path = os.path.join(script_dir, 'Python to PPT', 'Outputs', f'Template_chart_{industry}_{country_str}.pptx')
    ppttc_path = os.path.join(script_dir, 'Python to PPT', 'inputs', 'JSON.ppttc')
    template_path = os.path.join(script_dir, 'Python to PPT', 'Template files', "input.pptx")
    file_path_1 = os.path.join(script_dir, 'Python to PPT', 'Outputs', f'{industry}_{country_str}_PI_Data.pptx')
    file_path_2 = os.path.join(script_dir, 'Python to PPT', 'Outputs', f'Template_chart_{industry}_{country_str}.pptx')

    # Load the PowerPoint presentation template
    presentation = Presentation(file_path)

    # Function to inspect slide shapes (for debugging purposes)
    def inspect_slide_shapes(file_path, slide_index):
        prs = Presentation(file_path)
        slide = prs.slides[slide_index]
    
        for i, shape in enumerate(slide.shapes):
            print(f"Shape {i}: Type={shape.shape_type}, Name={shape.name}")


    # Function to populate a table shape in the slide with data
    def populate_table_with_data(table, df, font_size):
        company_column_index = 0
        center_align_columns = [1, 4, 6, 7, 8, 9, 10]
        comma_separator_columns = [6, 7, 8, 9]

        for i in range(1, min(len(df) + 1, len(table.rows))):
            for j in range(min(len(df.columns), len(table.columns))):
                cell = table.cell(i, j)
                value = df.iloc[i - 1, j]
                if isinstance(value, str) and value.strip() == '':
                    cell_value = ''
                else:
                    if j == 1:
                        cell_value = round(float(value) / 1000, 1)
                    elif j == 4:
                        cell_value = round(float(value) * 100, 1)
                    elif j == 10:
                        cell_value = f"{value}x" if isinstance(value, float) else value
                    elif j in comma_separator_columns:
                        cell_value = f"{int(float(value)):,}"
                    else:
                        cell_value = value

                cell.text = str(cell_value)

                # Apply direct paragraph formatting
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.space_before = Pt(0)
                    paragraph.space_after = Pt(0)
                    paragraph.line_spacing = 1.0
                    paragraph.level = 0

                    # Set alignment for specific columns
                    if j in center_align_columns:
                        paragraph.alignment = PP_ALIGN.CENTER
                    else:
                        paragraph.alignment = PP_ALIGN.LEFT

                    for run in paragraph.runs:
                        run.font.size = font_size
                        run.font.color.rgb = RGBColor(0, 0, 0)
                        if j == 0:
                            run.font.bold = False

                    # Remove any list styling and bullets
                    pPr = paragraph._element.get_or_add_pPr()

                    # Remove bullet points
                    buNone = OxmlElement("a:buNone")
                    pPr.insert(0, buNone)

                    # Clear any indentation
                    ind = pPr.find(".//a:ind", namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                    if ind is not None:
                        pPr.remove(ind)

                    # Set indentation explicitly to zero
                    new_ind = OxmlElement('a:ind')
                    new_ind.set('left', '0')
                    new_ind.set('hanging', '0')
                    pPr.append(new_ind)

                # Clear any text_frame list styles
                lstStyle = cell.text_frame._element.find(".//a:lstStyle", namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                if lstStyle is not None:
                    cell.text_frame._element.remove(lstStyle)

            # New code to remove rows where "Company" column is empty
        for i in reversed(range(1, len(table.rows))):  # Iterate in reverse to safely remove rows
            company_cell = table.cell(i, company_column_index)
            if company_cell.text.strip() == '':
                row_element = table._tbl.tr_lst[i]  # Access the row element using the table's tr_lst
                table._tbl.remove(row_element)  # Remove the row element from the table

    relationship_types = ['Current', 'Recent', 'Greyspace','Whitespace']
    selected_columns = [
        'Company','Revenue (M)','Country', 'Primary Industry', 'EBIT margin Change (% pts.)  [Benchmark year-Base year]',
        'Relative-EBIT Margin','Static EBIT opportunity (Median)', 'Static SG&A Opportunity (Median)', 
        'Static NWC opportunity (Median)','Median Incremental Rev.- EBIT Oppt. (Static Historical)',
        "Net debt/EBITDA (Dec'22)", 'Relative ESG Combined score'
    ]
    # chart_columns =['Company','Combined Opp. Score', 'Bain Relationship Score', 'Revenue (M)']
    # # df_chart_data = df_filtered[chart_columns]
    # data = df_chart_data.set_index('Company').rename(columns={
    #     'Combined Opp. Score': 'x',
    #     'Bain Relationship Score': 'y',
    #     'Revenue (M)': 'size'
    # }).to_dict('index')

    # Function to filter the dataset based on various criteria
    def filter_data_1(df, relationship):
        if relationship:
            if relationship == 'Current':
                df = df[df['Bain Relationship'].str.startswith('Current')]
            elif relationship == 'Recent':
                df = df[df['Bain Relationship'].str.startswith('Recent')]
            elif relationship == 'Greyspace':
                df = df[df['Bain Relationship'].str.startswith('Greyspace')]
            elif relationship == 'Whitespace':
                df = df[df['Bain Relationship'].str.startswith('Whitespace')]
        return df
    
    slide_index=0
    slide = presentation.slides[slide_index]
    title_shape = slide.shapes.title
    revenue_str = '<$1B' if revenue == 'less_than_100000' else '>between $1B and $5B' if revenue == 'between 10000_and_50000' else '>between $5B and $10B' if revenue == 'between 50000_and_100000' else 'between $10B and $20B' if revenue == 'between 100000_and_200000' else '>$20B' if revenue == 'more_than_200000' else "All"
    title_shape.text = f" PI Data for: {industry}, Revenue: {revenue_str} {priority_status}"


    slide_index = 2
    for relationship in relationship_types:
        filtered_df = filter_data_1(df_filtered, relationship)
        filtered_df = filtered_df[selected_columns]
        # print(filtered_df)
        if not filtered_df.empty:
            num_slides = (len(filtered_df) + 9) // 10
            for slide_num in range(num_slides):
                df_slice = filtered_df[slide_num * 10:(slide_num + 1) * 10]
                slide = presentation.slides[slide_index]
                title_shape = slide.shapes.title
                # print(revenue)
                revenue_str = '<$1B' if revenue == 'less_than_100000' else '>between $1B and $5B' if revenue == 'between 10000_and_50000' else '>between $5B and $10B' if revenue == 'between 50000_and_100000' else 'between $10B and $20B' if revenue == 'between 100000_and_200000' else '>$20B' if revenue == 'more_than_200000' else "All"
                title_shape.text = f"{industry} ({country_str}) Revenue: {revenue_str} {priority_status} ({relationship}) - ({slide_num + 1}/{num_slides})"

                table = next((shape.table for shape in slide.shapes if shape.has_table), None)
                if table:
                    populate_table_with_data(table, df_slice, font_size=Pt(10))

                # print(revenue_str)
                slide_index += 1
# Remove any remaining empty slides in the presentation
    while slide_index < len(presentation.slides):
        slide = presentation.slides[slide_index]
        xml_slides = presentation.slides._sldIdLst
        slide_id = slide.slide_id
        for s in xml_slides:
            if s.attrib['id'] == str(slide_id):
                xml_slides.remove(s)
                break
                
    # print(slide_index)


    # After removing the empty slides, create a new slide for the bubble chart
    slide2 = presentation.slides[1]

    # Create a copy of df_filtered and keep only the necessary columns
    ddf = df_filtered.copy()
    columns_to_keep = ['Company', 'Combined Opp. Score', 'Bain Relationship Score', 'Revenue (M)','Priority_status']
    # print(ddf[columns_to_keep])

    # Check if 'Bain Relationship Score' exists and filter the DataFrame
    if 'Bain Relationship Score' in ddf.columns:
        ddf = ddf[columns_to_keep]
        ddf.to_csv("filtered_csv.csv", index=False)  # Optional: save filtered data as CSV for reference
    else:
        print("Warning: 'Bain Relationship Score' column not found in the DataFrame")
        return  # Exit the function if required column is missing
    
    # Extract the relevant data for the bubble chart
    companies = ddf['Company'].tolist()  # Get company names
    x_values = ddf['Combined Opp. Score'].tolist()  # X-axis values
    y_values = ddf['Bain Relationship Score'].tolist()  # Y-axis values
    revenue = ddf['Revenue (M)'].tolist()  # Bubble sizes
    priority_comps = [True if status == 'Priority' else False for status in ddf['Priority_status'].tolist()]

    # Ensure that the values are numeric and handle any non-numeric data gracefully
    def safe_float_conversion(val):
        try:
            return float(val)
        except (ValueError, TypeError):
            return np.nan  # Return NaN for non-convertible values

    # Convert values to floats, handling non-numeric data
    x_values = [safe_float_conversion(x) for x in x_values]
    y_values = [safe_float_conversion(y) for y in y_values]
    revenue = [safe_float_conversion(r) for r in revenue]

    # Filter out data points where any value is NaN
    filtered_data = [
        (company, priority_comps, x, y, r)
        for company, priority_comps, x, y, r in zip(companies, priority_comps, x_values, y_values, revenue)
        if not (np.isnan(x) or np.isnan(y) or np.isnan(r))
    ]

    # Sort filtered_data by revenue (r, the 5th element in the tuple) in descending order
    filtered_data_sorted = sorted(filtered_data, key=lambda item: item[4], reverse=True)

        # Convert the sorted data to a DataFrame
    df = pd.DataFrame(filtered_data_sorted, columns=['Company', 'Priority Status', 'X Value', 'Y Value', 'Revenue'])

    # Write the DataFrame to an Excel file
    output_file = 'filtered_data_sorted.xlsx'
    df.to_excel(output_file, index=False)

    print(f"Data has been successfully written to {output_file}")



    # Define chart position and size BEFORE adding the chart
    x, y, cx, cy = Inches(1.5), Inches(1), Inches(10), Inches(5)

    # Create Bubble chart data
    chart_data = BubbleChartData()
    series = chart_data.add_series('Revenue Data')  # Set the series name

    # Add data points (x, y, bubble size) and store additional metadata
    for company, priority_comps, x_val, y_val, r_val in filtered_data_sorted:
        point = series.add_data_point(x_val, y_val, r_val)
        # Attach metadata to the point object manually for later use
        point.company = company
        point.priority = priority_comps

    # Add chart to slide2
    chart_shape = slide2.shapes.add_chart(
        XL_CHART_TYPE.BUBBLE, x, y, cx, cy, chart_data
    )
    chart = chart_shape.chart  # Get the chart object for further manipulation if necessary

    # Set color for specific bubbles based on the flag (priority status)
    for i, point in enumerate(chart.series[0].points):
        if filtered_data_sorted[i][1]:  # Check if the point is marked as priority
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = RGBColor(192, 0, 0)  # Red color for priority companies
        else:
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = RGBColor(100, 100, 100)  # Gray color for non-priority companies
            
        point.format.fill.transparency = 0.3  # Set transparency to 30% for all points

    # # Set color for specific bubbles based on the priority status
    # for i, point in enumerate(chart.series[0].points):
    #     # Check if the company is a priority company
    #     if point.priority[i] == 'Priority':
    #         point.format.fill.solid()
    #         point.format.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red color for priority companies
    #         point.format.fill.transparency = 0.3  # Set transparency to 30%
    #     else:
    #         point.format.fill.solid()
    #         point.format.fill.fore_color.rgb = RGBColor(100, 100, 100)  # Gray color for non-priority companies
    #         point.format.fill.transparency = 0.3  # Set transparency to 30%
        

    # Modify chart axis properties
    value_axis = chart.value_axis
    category_axis = chart.category_axis

    # Axis scaling
    value_axis.minimum_scale = 0
    value_axis.maximum_scale = 12  # Fix y-axis to 8
    category_axis.minimum_scale = 0
    category_axis.maximum_scale = 25  # Fix x-axis to 20

    # Reduce font size for x-axis (category axis)
    category_axis.tick_labels.font.size = Pt(10)  # Set x-axis font size to 10 pt

    # Reduce font size for y-axis (value axis)
    value_axis.tick_labels.font.size = Pt(10)  # Set y-axis font size to 10 pt

    # Remove all gridlines and tick marks except x and y axes
    value_axis.has_major_gridlines = False  # Remove horizontal gridlines
    category_axis.has_major_gridlines = False  # Remove vertical gridlines
    category_axis.major_tick_mark = XL_TICK_MARK.NONE
    value_axis.major_tick_mark = XL_TICK_MARK.NONE

    # Set chart axis titles (hide)
    category_axis.has_title = False
    value_axis.has_title = False
    chart.has_legend = False

    # Add static images for x and y axis labels
    script_dir = os.path.dirname(os.path.abspath(__file__))

    x_label_path = os.path.join(script_dir, 'static', 'bottom axis.png')
    y_label_path = os.path.join(script_dir, 'static', 'a.png')

    # Configure dimensions for x and y label images
    x_image_width = Inches(10.0)  # Adjust the width of the x-axis image
    x_image_height = Inches(0.5)  # Adjust the height of the x-axis image
    y_image_width = Inches(0.5)  # Adjust the width of the y-axis image
    y_image_height = Inches(4.75)  # Adjust the height of the y-axis image

    # Add x-axis label image (position near bottom of the chart)
    slide2.shapes.add_picture(x_label_path, Inches(1.5), Inches(6), width=x_image_width, height=x_image_height)

    # Add y-axis label image (position near the left of the chart, rotated)
    y_label_shape = slide2.shapes.add_picture(y_label_path, Inches(0.75), Inches(1.5), width=y_image_width, height=y_image_height)
    y_label_shape.rotation = 0  # Optionally rotate the y-axis label image

    # Dynamically set the title of the slide
    title_shape = slide2.shapes.title
    # revenue_str = '<$5B' if revenue == 'less_than_500000' else '>$5B' if revenue == 'greater_than_500000' else "All"

    # Assuming you have variables for industry, country_str, priority_status, relationship, slide_num, and num_slides
    title_shape.text = f"{industry} ({country_str}) Revenue: {revenue_str} {priority_status}"

    # Identify the top 5 companies by revenue
    top_5_companies = [filtered_data_sorted[i][0] for i in range(min(5, len(filtered_data_sorted)))]

    # Loop through all shapes in the slide
    for shape in slide2.shapes:
        if shape.has_chart:
            print("chart found")
            # Check if the chart's title matches 'Chart' (or whatever title you're checking)
            print(shape.chart.chart_title.text_frame.text)
            
            # Access the first series of the chart
            series = shape.chart.plots[0].series[0]
            series.has_data_labels = True  # Enable data labels for the series

            # Loop over all points in the series and show data labels for priority or top 5 companies
            for i, point in enumerate(series.points):
                company_name = filtered_data_sorted[i][0]  # Get the company name
                is_priority = filtered_data_sorted[i][1]  # Check if the company is tagged as priority

                # Show label if the company is in the top 5 by revenue or tagged as priority
                if company_name in top_5_companies or is_priority:
                    point.data_label.show_category_name = True  # Show category name
                    point.data_label.text_frame.text = company_name  # Set the company name as the label

                    # Set the font size for the data label
                    for paragraph in point.data_label.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(8)  # Set the font size to 8 points
                else:
                    # Hide the data label for companies not in top 5 or not tagged as priority
                    point.data_label.show_category_name = False
        

    # Save the presentation with both tables and the bubble chart
    presentation.save(file_path_1)


    return file_path_1

            

def load_data():
    global df, priority_companies_df, priority_companies
    
    filename1 = 'PI data.xlsx'
    filename2 = 'Client mapping.xlsx'
    sheet_name = 'Base year_2021'
    usecols = 'B:AL'
    
    script_dir = os.path.dirname(os.path.abspath(__file__))
    input_path = os.path.join(script_dir, app.config['INPUT_FOLDER'], filename1)
    mapping_path = os.path.join(script_dir, app.config['INPUT_FOLDER'], filename2)
    
    df = pd.read_excel(input_path, sheet_name=sheet_name, usecols=usecols)
    priority_companies_df = pd.read_excel(mapping_path)
    priority_companies = priority_companies_df['Informal client name'].tolist()
    
    # Add a column to indicate priority status based on the list of priority companies
    df['Priority_status'] = np.where(df['Company'].isin(priority_companies), 'Priority', 'Non-Priority')
    # Replace certain geographic values and fill NaN values with 0
    df['Geography'] = df['Geography'].replace('N. America', 'AMER')
    df = df.fillna(0)

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload_file', methods=['GET', 'POST'])
def upload_file():
    if request.method == 'POST':
        access_key = request.form.get('access_key')

        # Check if access key is correct
        if access_key != ACCESS_KEY:
            flash('Invalid access key. Please try again.')
            return redirect(request.url)

        file1 = request.files.get('file1')
        file2 = request.files.get('file2')

        if not file1 and not file2:
            flash('No files selected for upload.')
            return redirect(request.url)

        def save_file(file):
            if file and allowed_file(file.filename):
                filename = file.filename  # Directly use the original filename
                file_path = os.path.join(app.config['INPUT_FOLDER'], filename)
                
                # Check if file exists and delete it
                if os.path.exists(file_path):
                    os.remove(file_path)
                    flash(f'{filename} was replaced.')

                # Save the new file
                file.save(file_path)
                flash(f'{filename} uploaded successfully!')

        # Handle the first file
        if file1:
            save_file(file1)

        # Handle the second file
        if file2:
            save_file(file2)

        # Reload data after files are uploaded
        load_data()

        # After processing the files, redirect to the index page
        return redirect(url_for('index'))

    return render_template('upload_file.html', correct_access_key=ACCESS_KEY)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in {'xls', 'xlsx'}
############################################################################################
if __name__ == '__main__':
    app.run(debug=True)
