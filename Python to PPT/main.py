#
import pandas as pd
import os
import json
import subprocess
import csv
import time
from collections import defaultdict
from pptx import Presentation
from pptx.util import Pt
import copy
import win32com.client

script_dir = os.path.dirname(os.path.abspath(__file__)) 

# Function to read an Excel file and store data in a DataFrame
def read_excel_to_dataframe(file_path, sheet_name, start_row=0, start_col=0):
    # Read the Excel file, specifying the sheet name, start row, and column
    df = pd.read_excel(file_path, sheet_name=sheet_name, header=start_row)
 
   # Select columns starting from start_col
    df = df.iloc[:, start_col:]
 
    # Display the first few rows of the DataFrame to verify
    #print("Headers:", df.columns.tolist())
    #print(df.head())
   
    return df
 
# Function to filter the DataFrame based on specified columns
def filter_dataframe(df, industry=None, country=None, revenue=None,Relationship =None):
    if industry is not None:
        df = df[df['Bain Industry'] == industry]
    if country is not None:
        df = df[df['Country'] == country]
    if Relationship is not None:
        df = df[df['Bain Relationship'].str.contains(Relationship,case =True, na=False)]
    if revenue is not None:
        if revenue_filter_type == 'greater':
            df = df[df['Revenue (M)'] > revenue]
        elif revenue_filter_type == 'less':
            df = df[df['Revenue (M)'] < revenue]
 
    # Display the filtered DataFrame
    #print("Filtered DataFrame:")
    #print(df.head())
    # Print the number of rows in the filtered DataFrame
    #print(f"Number of rows in filtered DataFrame: {df.shape[0]}")
   
    #selected_columns = ['Company','Revenue (M)','Primary Industry','EBIT margin Change (% pts.)  [Benchmark year-Base year]','Static EBIT opportunity (Median)','Static SG&A Opportunity (Median)','Static NWC opportunity (Median)','Median Incremental Rev. - EBIT Oppty (Static Projected)','Historical rev. CAGR',"Net debt/EBITDA (Dec'22)",'Relative ESG Combined score']
    #new_df =df[selected_columns]
 
    #print(new_df.head())
    return df
 
 
# Replace 'your_file.xlsx' with the path to your actual Excel file
file_path = os.path.join(script_dir, 'inputs', '20230411_A1OI_Summary Update_v4.xlsx')
 
# Specify the starting row and column
sheet_name = 'Base year_2021'  # Example: Name of the sheet to read from
start_row = 7  # Example: Start reading from the second row (index 1)
start_col = 1  # Example: Start reading from the second column (index 1)
 
# Read the Excel file and store the data in a DataFrame
dataframe = read_excel_to_dataframe(file_path, sheet_name, start_row=start_row, start_col=start_col)

# Specify the filter criteria
industry_filter = 'Media'  # Example: Filter by 'Technology' industry
country_filter = 'France'  # Example: Filter by 'USA' country
Relationship_filter = 'Greyspace'
revenue_filter = 5000000  # Example: Filter by revenue greater than 1000000
revenue_filter_type = 'less'  # 'greater' for >, 'less' for <
 
filtered_dataframe = filter_dataframe(dataframe, industry=industry_filter, country=country_filter, revenue=revenue_filter,Relationship=Relationship_filter)
#print(filtered_dataframe)
filtered_dataframe_1 = filter_dataframe(dataframe, industry=industry_filter, country=country_filter, revenue=revenue_filter)

# filtered_dataframe.to_csv(industry_filter+".csv")

 
# Path to your existing PowerPoint file
file_path = os.path.join(script_dir, 'inputs', 'Template slides for Transformation Screens.pptx')

file_path_1 = os.path.join(script_dir, 'Outputs', 'Template_tables_{}.pptx'.format(industry_filter))
 
# Open the PowerPoint file
presentation = Presentation(file_path)
 

 
# Define the list with the specified values
keywords = ['Current', 'Recent', 'Greyspace']
 
 
def update_table_in_slide(pres, slide_index, dataframe, title_text):
    # Get the slide to be updated
    slide = pres.slides[slide_index]
    #new_slide_layout = slide.slide_layout    
    #new_slide = presentation.slides.add_slide(new_slide_layout)
 
    # Update the title of the slide
    title = slide.shapes.title
    title.text = title_text
   
    # Find the table in the slide
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
           
            # Update the table with data from the dataframe
            for row_idx in range(len(dataframe)):
                for col_idx in range(len(dataframe.columns)):
                    cell = table.cell(row_idx + 1, col_idx)  # +1 to skip the header row
                    cell.text = str(dataframe.iat[row_idx, col_idx])
 
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(10)
                            run.font.name = 'Arial'
            return True
       
    return False
 
for slide_index in range(1, 4):
    #print(keywords[slide_index-1])
    filtered_dataframe = filter_dataframe(dataframe, industry=industry_filter, country=country_filter, revenue=revenue_filter, Relationship=keywords[slide_index-1])
    
    # filtered_dataframe.to_csv(industry_filter+".csv") 
    title_text = f"{industry_filter} ({country_filter}) Revenue: {revenue_filter_type} than {revenue_filter} {keywords[slide_index-1]}"

    # Create a new dataframe using the function's output
    df_new = pd.DataFrame(filtered_dataframe)
    selected_columns = ['Company', 'Revenue (M)', 'Primary Industry', 'EBIT margin Change (% pts.)  [Benchmark year-Base year]', 'Static EBIT opportunity (Median)', 'Static SG&A Opportunity (Median)', 'Static NWC opportunity (Median)', 'Median Incremental Rev. - EBIT Oppty (Static Projected)', 'Historical rev. CAGR', "Net debt/EBITDA (Dec'22)", 'Relative ESG Combined score']
    new_df = df_new[selected_columns].head(5)
    #new_df.to_csv(industry_filter+".csv")
    update_table_in_slide(presentation, slide_index, new_df, title_text)
    presentation.save(file_path_1)
    #print(f"Presentation saved as '{file_path}'")

ddf = filtered_dataframe_1.copy()
columns_to_keep = ['Company', 'Combined Opp. Score', 'Bain Relationship Score  ', 'Revenue (M)']
ddf = ddf[columns_to_keep]
ddf.to_csv("filtered_csv.csv", index = False)
# Save the modified presentation back to the same file

 


def generate_json_for_bubble_chart(data, template_path):
    # Convert data to a defaultdict for easier access.
    data = defaultdict(lambda: {'x': 0, 'y': 0, 'size': 0}, data)

    # Sort companies.
    companies = sorted(data.keys())

    # Set up the basic structure of the JSON object.
    chart_data = {
        "template": template_path,
        "data": [
            {
                "name": "BubbleChart1",
                "table": [
                    # Header row with column labels.
                    [{"string": "Company"}, {"string": "Combined Opp. Score"}, {"string": "Bain Relationship Score  "}, {"string": "Revenue (M)"}],
                    # Rows for each data point.
                    *[
                        [{"string": company}, {"number": data[company]['x']}, {"number": data[company]['y']}, {"number": data[company]['size']}]
                        for company in companies
                    ]
                ]
            }
        ]
    }

    return json.dumps([chart_data], indent=4)

def run_thinkcell_cli(ppttc_file, output_pptx_path):
    command = [
        "C:\\Program Files (x86)\\think-cell\\ppttc.exe", ppttc_file, '-o', output_pptx_path
    ]

    # Execute the command using subprocess.
    try:
        result = subprocess.run(
            command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
        print("think-cell processing successful.")
        print("Command output:", result.stdout)
    except subprocess.CalledProcessError as e:
        print("Error:", e)
        print("Standard Output:", e.stdout)
        print("Standard Error:", e.stderr)

def read_csv(data_file_path):
    data = {}
    with open(data_file_path, 'r', encoding='utf-8-sig') as file:
        reader = csv.DictReader(file)
        headers = reader.fieldnames
        if not all(col in headers for col in ['Company', 'Combined Opp. Score', 'Bain Relationship Score  ', 'Revenue (M)']):
            raise ValueError(f"CSV file must contain 'Company', 'Combined Opp. Score', 'Bain Relationship Score  ', and 'Revenue (M)' columns. Found headers: {headers}")

        for row in reader:
            company = row['Company']
            data[company] = {
                'x': float(row['Combined Opp. Score']),
                'y': float(row['Bain Relationship Score  ']),
                'size': float(row['Revenue (M)'])
            }

    return data

# Paths
#file_new_path ="Template_chart.pptx"
template_path = os.path.join(script_dir, 'Template files', "input.pptx")
os.path.join(script_dir, 'inputs','JSON.ppttc')
ppttc_path = os.path.join(script_dir, 'inputs','JSON.ppttc')
#output_pptx_path = r"C:\Users\62320\OneDrive - Bain\2024-Viz\Python Web App\Python to PPT\Outputs\Template_chart_Media_France.pptx"
output_pptx_path = os.path.join(script_dir, 'Outputs', 'Template_chart_{}_{}.pptx'.format(industry_filter, country_filter))
data_file_path = "filtered_csv.csv"
#slide_index = 0

# Delete old files if they exist
if os.path.exists(ppttc_path):
    os.remove(ppttc_path)


# Read data from csv
data = read_csv(data_file_path)

# Generate JSON for bubble chart
json_output = generate_json_for_bubble_chart(data, template_path)
with open(ppttc_path, 'w') as file:
    file.write(json_output)

# Debug: Print the generated JSON
#print("Generated JSON:")
#print(json_output)

# Check if .ppttc file is updated
if os.path.exists(ppttc_path):
    print(f"{ppttc_path} was successfully created.")
    print("Last modified:", time.ctime(os.path.getmtime(ppttc_path)))
else:
    print(f"Error: {ppttc_path} was not created.")

# Verify contents of the .ppttc file
with open(ppttc_path, 'r') as file:
    content = file.read()
    #print("Contents of the .ppttc file:")
    #print(content)

# Run the think-cell CLI.
run_thinkcell_cli(ppttc_path, output_pptx_path)

# Check if output .pptx file is created
if os.path.exists(output_pptx_path):
    print(f"{output_pptx_path} was successfully created.")
    print("Last modified:", time.ctime(os.path.getmtime(output_pptx_path)))
else:
    print(f"Error: {output_pptx_path} was not created.")

############################################################
#Slide - Title

press = Presentation(output_pptx_path)
first_slide = press.slides[0]

# Access the title shape of the first slide
title_shape = first_slide.shapes.title

title_text_bubble = f"{industry_filter} ({country_filter}) Revenue: {revenue_filter_type} than {revenue_filter}"
# Change the title text
title_shape.text = title_text_bubble

# Save the presentation
press.save(output_pptx_path)

######################################################################


lst=[file_path_1,output_pptx_path]
out_path = os.path.join(script_dir, 'Outputs', 'Combined_Template_chart_{}_{}.pptx'.format(industry_filter, country_filter))
def merge_presentations(presentations, path):
  ppt_instance = win32com.client.Dispatch('PowerPoint.Application')
  prs = ppt_instance.Presentations.open(os.path.abspath(presentations[0]), True, False, False)

  for i in range(1, len(presentations)):
      prs.Slides.InsertFromFile(os.path.abspath(presentations[i]), prs.Slides.Count)

  prs.SaveAs(os.path.abspath(path))
  prs.Close()

merge_presentations(lst,out_path)

######################################################################

def delete_pptx_files(files):
    for file in files:
        if os.path.exists(file) and file.endswith('.pptx'):
            os.remove(file)
        else:
            print("")

# List of .pptx files to be deleted
files_to_delete = [file_path_1,output_pptx_path]

delete_pptx_files(files_to_delete)

