from flask import Flask, render_template, request, send_file
import pandas as pd
import os
import csv
from collections import defaultdict
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
import numpy as np
from io import BytesIO
import xml.etree.ElementTree as ET
import json  # Import the json module
from collections import defaultdict
import subprocess

app = Flask(__name__)

# Load data once when the app starts
sheet_name = 'Base year_2021'
usecols = 'B:AK'
df = pd.read_excel('20230411_A1OI_Summary Update_v4.xlsx', sheet_name=sheet_name, usecols=usecols)
priority_companies_df = pd.read_excel('A1OI Priority client mapping.xlsx')
priority_companies = priority_companies_df['Informal client name'].tolist()
df['Priority_status'] = np.where(df['Company'].isin(priority_companies_df['Informal client name']), 'Priority', 'Non-Priority')
df['Geography'] = df['Geography'].replace('N. America', 'AMER')
df = df.fillna(0)

def get_contextual_values(df, geography=None, industry=None):
    if geography:
        filtered_df = df[df['Geography'] == geography]
    else:
        filtered_df = df
    
    if industry:
        filtered_df = filtered_df[filtered_df['Bain Industry'] == industry]
    
    geographies = sorted(df['Geography'].unique().tolist())
    countries = sorted(filtered_df['Country'].unique().tolist()) if geography else sorted(df['Country'].unique().tolist())
    industries = sorted(df['Bain Industry'].unique().tolist())
    industries_primary = sorted(filtered_df['Primary Industry'].unique().tolist())
    relationships = sorted(df['Bain Relationship'].unique().tolist())
    priority_column = sorted(df['Priority_status'].unique().tolist())

    return industries, relationships, countries, geographies, industries_primary, priority_column

def filter_data(df, industry, relationship, revenue, geography, country, industry_primary, combined_opp_score, priority_status):
    df['Country'] = df['Country'].str.strip()

    if isinstance(country, str):
        country = [c.strip() for c in country.split(',')]
    else:
        country = [c.strip() for c in country]

    if industry:
        df = df[df['Bain Industry'] == industry]
    if industry_primary:
        df = df[df['Primary Industry'] == industry_primary]
    if relationship:
        if relationship == 'Current':
            df = df[df['Bain Relationship'].str.startswith('Current')]
        elif relationship == 'Recent':
            df = df[df['Bain Relationship'].str.startswith('Recent')]
        elif relationship == 'Greyspace':
            df = df[df['Bain Relationship'].str.startswith('Greyspace')]
    if revenue == 'less_than_500000':
        df = df[df['Revenue (M)'] < 5000]
    elif revenue == 'greater_than_500000':
        df = df[df['Revenue (M)'] > 5000]
    if priority_status:
        df = df[df['Priority_status'] == priority_status]
    if geography:
        df = df[df['Geography'] == geography]
    if country:
        df = df[df['Country'].isin(country)]
    if combined_opp_score:
        if combined_opp_score == 'less_than_5':
            df = df[df['Combined Opp. Score'] < 5]
        elif combined_opp_score == 'between_5_and_10':
            df = df[(df['Combined Opp. Score'] >= 5) & (df['Combined Opp. Score'] <= 10)]
        elif combined_opp_score == 'greater_than_10':
            df = df[df['Combined Opp. Score'] > 10]
    
    return df

@app.route('/')
def home():
    industry = request.args.get('industry')
    relationship = request.args.get('relationship')
    revenue = request.args.get('revenue')
    country = request.args.getlist('country[]')
    geography = request.args.get('geography')
    industry_primary = request.args.get('industry_primary')
    combined_opp_score = request.args.get('combined_opp_score')
    priority_status = request.args.get('priority_status')

    filters_applied = any([industry, relationship, revenue, country, geography, industry_primary, combined_opp_score, priority_status])

    df_filtered = filter_data(df, industry, relationship, revenue, geography, country, industry_primary, combined_opp_score, priority_status)
    filtered_data = df_filtered.to_dict(orient='records')
    min_revenue = df_filtered['Revenue (M)'].min()

    industries, relationships, countries, geographies, industries_primary, priority_column = get_contextual_values(df, geography, industry)

    return render_template('index.html',
                           industries=industries,
                           industries_primary=industries_primary,
                           relationships=relationships,
                           geographies=geographies,
                           countries=countries,
                           revenue_ranges=['less_than_500000', 'greater_than_500000'],
                           filtered_data=filtered_data,
                           filters_applied=filters_applied,
                           industry=industry,
                           industry_primary=industry_primary,
                           revenue=revenue,
                           relationship=relationship,
                           geography=geography,
                           country=country,
                           min_revenue=min_revenue,
                           priority_companies=priority_companies,
                           priority_column=priority_column,
                           priority_status=priority_status)

@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    industry = request.form.get('industry')
    relationship = request.form.get('relationship')
    revenue = request.form.get('revenue')
    country = request.form.get('country[]', '').split('|')
    country = [c.strip() for c in country if c.strip()]
    geography = request.form.get('geography')
    industry_primary = request.form.get('industry_primary')
    combined_opp_score = request.form.get('combined_opp_score')
    priority_status = request.form.get('priority_status')

    df_filtered = filter_data(df, industry, relationship, revenue, geography, country, industry_primary, combined_opp_score, priority_status)

    ppt_path = generate_ppt_file(df_filtered, industry, country, revenue, geography, industry_primary, combined_opp_score, priority_status)

    return send_file(ppt_path, as_attachment=True)

def generate_ppt_file(df_filtered, industry, country, revenue, geography, industry_primary, combined_opp_score, priority_status):
    script_dir = os.path.dirname(os.path.abspath(__file__))
    file_path = os.path.join(script_dir, 'Python to PPT', 'inputs', 'Template_v1.pptx')
    country_str = ', '.join(country) if country else 'No Country Selected'
    file_path_1 = os.path.join(script_dir, 'Python to PPT', 'Outputs', f'Template_tables_{industry}_{country_str}.pptx')

    presentation = Presentation(file_path)

    def inspect_slide_shapes(file_path, slide_index):
        prs = Presentation(file_path)
        slide = prs.slides[slide_index]
    
        for i, shape in enumerate(slide.shapes):
            print(f"Shape {i}: Type={shape.shape_type}, Name={shape.name}")



    def populate_table_with_data(table, df, font_size):
        center_align_columns = [1, 4, 6, 7, 8, 9, 10]
        comma_separator_columns = [6, 7, 8, 9]

        for i in range(1, min(len(df) + 1, len(table.rows))):
            for j in range(min(len(df.columns), len(table.columns))):
                cell = table.cell(i, j)
                value = df.iloc[i - 1, j]
                if isinstance(value, str) and value.strip() == '':
                    cell_value = ''
                else:
                    if j == 1:
                        cell_value = round(float(value) / 1000, 1)
                    elif j == 4:
                        cell_value = round(float(value) * 100, 1)
                    elif j == 10:
                        cell_value = f"{value}x" if isinstance(value, float) else value
                    elif j in comma_separator_columns:
                        cell_value = f"{int(float(value)):,}"
                    else:
                        cell_value = value

                cell.text = str(cell_value)

                # Apply direct paragraph formatting
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.space_before = Pt(0)
                    paragraph.space_after = Pt(0)
                    paragraph.line_spacing = 1.0
                    paragraph.level = 0

                    # Set alignment for specific columns
                    if j in center_align_columns:
                        paragraph.alignment = PP_ALIGN.CENTER
                    else:
                        paragraph.alignment = PP_ALIGN.LEFT

                    for run in paragraph.runs:
                        run.font.size = font_size
                        run.font.color.rgb = RGBColor(0, 0, 0)
                        if j == 0:
                            run.font.bold = False

                    # Remove any list styling and bullets
                    pPr = paragraph._element.get_or_add_pPr()

                    # Remove bullet points
                    buNone = OxmlElement("a:buNone")
                    pPr.insert(0, buNone)

                    # Clear any indentation
                    ind = pPr.find(".//a:ind", namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                    if ind is not None:
                        pPr.remove(ind)

                    # Set indentation explicitly to zero
                    new_ind = OxmlElement('a:ind')
                    new_ind.set('left', '0')
                    new_ind.set('hanging', '0')
                    pPr.append(new_ind)

                # Clear any text_frame list styles
                lstStyle = cell.text_frame._element.find(".//a:lstStyle", namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                if lstStyle is not None:
                    cell.text_frame._element.remove(lstStyle)

    relationship_types = ['Current', 'Recent', 'Greyspace','Whitespace']
    selected_columns = [
        'Company','Revenue (M)','Country', 'Primary Industry', 'EBIT margin Change (% pts.)  [Benchmark year-Base year]',
        'Relative-EBIT Margin','Static EBIT opportunity (Median)', 'Static SG&A Opportunity (Median)', 
        'Static NWC opportunity (Median)','Median Incremental Rev.- EBIT Oppt. (Static Historical)',
        "Net debt/EBITDA (Dec'22)", 'Relative ESG Combined score'
    ]
    chart_columns =['Company','Combined Opp. Score', 'Bain Relationship Score', 'Revenue (M)']
    df_chart_data = df_filtered[chart_columns]
    data = df_chart_data.set_index('Company').rename(columns={
        'Combined Opp. Score': 'x',
        'Bain Relationship Score': 'y',
        'Revenue (M)': 'size'
    }).to_dict('index')

        # Optionally, handle missing data with defaults
    data = defaultdict(lambda: {'x': 0, 'y': 0, 'size': 0}, data)
    ppttc_path = os.path.join(script_dir, 'Python to PPT', 'inputs', 'JSON.ppttc')
    def generate_json_for_bubble_chart(data, template_path,chart_name="Chart 146"):
        data = defaultdict(lambda: {'x': 0, 'y': 0, 'size': 0}, data)
        companies = sorted(data.keys())
        chart_data = {
            "template": template_path,
            "data": [
                {
                    "name": chart_name,
                    "table": [
                        [{"string": "Company"}, {"string": "Combined Opp. Score"}, {"string": "Bain Relationship Score"}, {"string": "Revenue (M)"}],
                        *[
                            [{"string": company}, {"number": data[company]['x']}, {"number": data[company]['y']}, {"number": data[company]['size']}]
                            for company in companies
                        ]
                    ]
                }
            ]
        }
        return json.dumps([chart_data], indent=4)
    def run_thinkcell_cli(ppttc_file, file_path):
        command = [
            "C:\\Program Files (x86)\\think-cell\\ppttc.exe", ppttc_file, '-o', file_path
        ]
        try:
            result = subprocess.run(
                command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
        except subprocess.CalledProcessError as e:
            print("Standard Error:", e.stderr)
    #inspect_slide_shapes(file_path,slide_index=1)
    
    json_output = generate_json_for_bubble_chart(data, file_path,chart_name="Chart 146")
    
    run_thinkcell_cli(ppttc_path, file_path)
    
    with open(ppttc_path, 'w') as f:
        f.write(json_output)
    
    slide_index =1
    #inspect_slide_shapes(file_path, slide_index)
    #update_think_cell_bubble_chart(file_path,slide_index,7,None)

    slide_index = 2
    for relationship in relationship_types:
        filtered_df = filter_data(df_filtered, industry, relationship, revenue, geography, country, industry_primary, combined_opp_score, priority_status)
        filtered_df = filtered_df[selected_columns]
        print(filtered_df)
        if not filtered_df.empty:
            num_slides = (len(filtered_df) + 9) // 10
            for slide_num in range(num_slides):
                df_slice = filtered_df[slide_num * 10:(slide_num + 1) * 10]
                slide = presentation.slides[slide_index]
                title_shape = slide.shapes.title
                revenue_str = '<$5B' if revenue == 'less_than_500000' else '>$5B' if revenue == 'greater_than_500000' else "All"
                title_shape.text = f"{industry} ({country_str}) Revenue: {revenue_str} {priority_status} ({relationship}) - ({slide_num + 1}/{num_slides})"

                table = next((shape.table for shape in slide.shapes if shape.has_table), None)
                if table:
                    populate_table_with_data(table, df_slice, font_size=Pt(10))

                slide_index += 1
                



    while slide_index < len(presentation.slides):
        slide = presentation.slides[slide_index]
        xml_slides = presentation.slides._sldIdLst
        slide_id = slide.slide_id
        for s in xml_slides:
            if s.attrib['id'] == str(slide_id):
                xml_slides.remove(s)
                break

    presentation.save(file_path_1)

    return file_path_1  # Return the path to the presentation with only tables

# Commenting out the code related to bubble chart generation and merging
# ddf = df_filtered.copy()
# columns_to_keep = ['Company', 'Combined Opp. Score', 'Bain Relationship Score', 'Revenue (M)']
# 
# if 'Bain Relationship Score' in ddf.columns:
#     ddf = ddf[columns_to_keep]
#     ddf.to_csv("filtered_csv.csv", index=False)
# else:
#     print("Warning: 'Bain Relationship Score' column not found in the DataFrame")
# 
# ppttc_path = os.path.join(script_dir, 'Python to PPT', 'inputs', 'JSON.ppttc')
# output_pptx_path = os.path.join(script_dir, 'Python to PPT', 'Outputs', f'Template_chart_{industry}_{country_str}.pptx')
# 
# if os.path.exists(ppttc_path):
#     os.remove(ppttc_path)
# 
# def generate_json_for_bubble_chart(data, template_path):
#     data = defaultdict(lambda: {'x': 0, 'y': 0, 'size': 0}, data)
#     companies = sorted(data.keys())
#     chart_data = {
#         "template": template_path,
#         "data": [
#             {
#                 "name": "BubbleChart1",
#                 "table": [
#                     [{"string": "Company"}, {"string": "Combined Opp. Score"}, {"string": "Bain Relationship Score"}, {"string": "Revenue (M)"}],
#                     *[
#                         [{"string": company}, {"number": data[company]['x']}, {"number": data[company]['y']}, {"number": data[company]['size']}]
#                         for company in companies
#                     ]
#                 ]
#             }
#         ]
#     }
#     return json.dumps([chart_data], indent=4)
# 
# def run_thinkcell_cli(ppttc_file, output_pptx_path):
#     command = [
#         "C:\\Program Files (x86)\\think-cell\\ppttc.exe", ppttc_file, '-o', output_pptx_path
#     ]
#     try:
#         result = subprocess.run(
#             command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
#     except subprocess.CalledProcessError as e:
#         print("Standard Error:", e.stderr)
# 
# def read_csv(data_file_path):
#     data = {}
#     with open(data_file_path, 'r', encoding='utf-8-sig') as file:
#         reader = csv.DictReader(file)
#         headers = reader.fieldnames
#         if not all(col in headers for col in ['Company', 'Combined Opp. Score', 'Bain Relationship Score', 'Revenue (M)']):
#             raise ValueError(f"CSV file must contain 'Company', 'Combined Opp. Score', 'Bain Relationship Score', and 'Revenue (M)' columns. Found headers: {headers}")
#         for row in reader:
#             company = row['Company']
#             data[company] = {
#                 'x': float(row['Combined Opp. Score']),
#                 'y': float(row['Bain Relationship Score']),
#                 'size': float(row['Revenue (M)'])
#             }
#     return data
# 
# data_file_path = "filtered_csv.csv"
# data = read_csv(data_file_path)
# template_path = os.path.join(script_dir, 'Python to PPT', 'Template files', "input.pptx")
# json_output = generate_json_for_bubble_chart(data, template_path)
# with open(ppttc_path, 'w') as file:
#     file.write(json_output)
# 
# run_thinkcell_cli(ppttc_path, output_pptx_path)
# 
# pythoncom.CoInitialize()
# try:
#     press = Presentation(output_pptx_path)
#     first_slide = press.slides[0]
#     title_shape = first_slide.shapes.title
#     title_text_bubble = f"{industry} ({country_str}) Revenue: {revenue}"
#     title_shape.text = title_text_bubble
#     press.save(output_pptx_path)
# finally:
#     pythoncom.CoUninitialize()
# 
# lst = [file_path_1, output_pptx_path]
# out_path = os.path.join(script_dir, 'Python to PPT', 'Outputs', f'Combined_Template_chart_{industry}_{country_str}.pptx')
# 
# def merge_presentations(presentations, path):
#     pythoncom.CoInitialize()
#     try:
#         ppt_instance = win32com.client.Dispatch('PowerPoint.Application')
#         prs = ppt_instance.Presentations.open(os.path.abspath(presentations[0]), True, False, False)
#         for i in range(1, len(presentations)):
#             prs.Slides.InsertFromFile(os.path.abspath(presentations[i]), prs.Slides.Count)
#         prs.SaveAs(os.path.abspath(path))
#         prs.Close()
#     finally:
#         pythoncom.CoUninitialize()
# 
# merge_presentations(lst, out_path)
# 
# def delete_pptx_files(files):
#     for file in files:
#         if os.path.exists(file) and file.endswith('.pptx'):
#             os.remove(file)
#         else:
#             print("")
# 
# files_to_delete = [file_path_1, output_pptx_path]
# delete_pptx_files(files_to_delete)

################################################
if __name__ == '__main__':
    app.run(debug=True)
